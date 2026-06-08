import os
import re
import unicodedata
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation


class LocalRecursiveTextSplitter:
    def __init__(self, chunk_size=1000, chunk_overlap=200):
        self.chunk_size = chunk_size
        self.chunk_overlap = min(chunk_overlap, max(0, chunk_size - 1))
        self.separators = ["\n\n", "\n", ". ", "! ", "? ", "; ", ", ", " "]

    def split_text(self, text):
        text = text.strip()
        if len(text) <= self.chunk_size:
            return [text] if text else []

        chunks = []
        start = 0
        while start < len(text):
            end = min(start + self.chunk_size, len(text))
            chunk_end = end

            if end < len(text):
                window = text[start:end]
                best_pos = -1
                for separator in self.separators:
                    pos = window.rfind(separator)
                    if pos > best_pos and pos > self.chunk_size * 0.45:
                        best_pos = pos + len(separator)
                if best_pos > 0:
                    chunk_end = start + best_pos

            chunk = text[start:chunk_end].strip()
            if chunk:
                chunks.append(chunk)

            if chunk_end >= len(text):
                break
            start = max(chunk_end - self.chunk_overlap, start + 1)

        return chunks


class DocumentProcessor:
    def __init__(self, chunk_size=1000, chunk_overlap=200):
        self.chunk_size = int(os.getenv("CHUNK_SIZE", chunk_size))
        self.chunk_overlap = int(os.getenv("CHUNK_OVERLAP", chunk_overlap))
        self.text_splitter = LocalRecursiveTextSplitter(self.chunk_size, self.chunk_overlap)

    def clean_text(self, text):
        text = re.sub(r"(\w)-\s*\n\s*(\w)", r"\1\2", text or "")
        text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text)
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        lines = [line.strip() for line in text.split("\n")]
        return "\n".join(lines).strip()

    def normalize_ascii(self, text):
        text = unicodedata.normalize("NFD", str(text or "").lower())
        text = "".join(ch for ch in text if unicodedata.category(ch) != "Mn")
        return re.sub(r"\s+", " ", text).strip()

    def parse_chapter_heading(self, line):
        clean = re.sub(r"\s+", " ", (line or "").strip())
        if not clean:
            return None

        patterns = [
            r"^(?:chapter|chuong)\s+([0-9]{1,2})\s*[:.\-]?\s*(.*)$",
            r"^([0-9]{1,2})\s+(Introduction|Overview of the UML Notation|Reliable, Scalable, and Maintainable Applications|Data Models and Query Languages)\b(.*)$",
        ]
        for pattern in patterns:
            match = re.match(pattern, clean, re.IGNORECASE)
            if not match:
                continue
            number = int(match.group(1))
            if not (1 <= number <= 40):
                return None
            title = ""
            if len(match.groups()) >= 2:
                title = re.sub(r"\s+", " ", " ".join(part for part in match.groups()[1:] if part)).strip(" .:-")
            title_norm = self.normalize_ascii(title)
            if title_norm and any(term in title_norm for term in [
                "described in section", "briefly described", "is described", "see chapter"
            ]):
                return None
            if title and not re.match(r"^[A-Za-z0-9]", title):
                return None
            return {"number": number, "title": title}
        return None

    def parse_section_heading(self, line):
        clean = re.sub(r"\s+", " ", (line or "").strip())
        match = re.match(r"^([0-9]{1,2}\.[0-9]+(?:\.[0-9]+)*)\s+(.{3,140})$", clean)
        if not match:
            return None
        title = match.group(2).strip(" .:-")
        if len(title) < 3:
            return None
        return {"number": match.group(1), "title": title}

    def detect_content_zone(self, text, page_number=None, heading=""):
        sample = self.normalize_ascii(f"{heading}\n{text[:1800]}")
        if any(term in sample for term in [
            "table of contents", "contents", "annotated table of contents", "muc luc"
        ]):
            return "toc"
        if any(term in sample for term in [
            "answers to exercises", "answer key", "solutions to exercises",
            "chapter 1 introduction 1 b", "multiple choice questions"
        ]):
            return "answer_key"
        if "references" in sample or "bibliography" in sample:
            return "references"
        if any(term in sample for term in ["preface", "acknowledgments", "acknowledgements"]):
            return "preface"
        if any(term in sample for term in ["appendix", "appendices"]):
            return "appendix"
        return "body"

    def make_section_path(self, chapter_number=0, chapter_title="", section_number="", section_title="", heading=""):
        parts = []
        if chapter_number:
            chapter_label = f"Chapter {chapter_number}"
            if chapter_title:
                chapter_label += f": {chapter_title}"
            parts.append(chapter_label)
        if section_number:
            section_label = section_number
            if section_title:
                section_label += f" {section_title}"
            parts.append(section_label)
        elif heading and heading not in parts:
            parts.append(heading)
        return " > ".join(part for part in parts if part)

    def detect_heading(self, text):
        lines = [line.strip() for line in self.clean_text(text).split("\n") if line.strip()]
        if not lines:
            return ""
        return self.detect_heading_line(lines[0])

    def detect_heading_line(self, line):
        line = re.sub(r"\s+", " ", (line or "").strip())[:180]
        if not line:
            return ""
        if re.match(r"^(chapter|chuong)\s+\d+(\s*[:.\-]\s*|\s+).+", line, re.IGNORECASE):
            return line
        if re.match(r"^(chapter|chuong)\s+\d+\s*$", line, re.IGNORECASE):
            return line
        if re.match(r"^\d+(\.\d+){0,4}\s+[A-Z][A-Za-z0-9 ,&:/()'\-]{3,120}$", line):
            return line
        if 8 <= len(line) <= 120 and line.isupper() and len(line.split()) <= 14:
            return line.title()
        return ""

    def detect_page_chapter(self, text):
        lines = [line.strip() for line in self.clean_text(text).split("\n") if line.strip()]
        for index, line in enumerate(lines[:14]):
            chapter = self.parse_chapter_heading(line)
            if chapter:
                if not chapter.get("title"):
                    for next_line in lines[index + 1:index + 4]:
                        if self.parse_chapter_heading(next_line):
                            continue
                        if 3 <= len(next_line) <= 140:
                            candidate = re.sub(r"\s+", " ", next_line).strip(" .:-")
                            candidate_norm = self.normalize_ascii(candidate)
                            if candidate and not re.match(r"^[A-Za-z0-9]", candidate):
                                return None
                            if any(term in candidate_norm for term in [
                                "described in section", "briefly described", "is described", "see chapter"
                            ]):
                                return None
                            chapter["title"] = candidate
                            break
                return chapter
        return None

    def split_text_into_units(self, text, page_number=None, slide_number=None, inherited_state=None):
        lines = [line.strip() for line in self.clean_text(text).split("\n")]
        units = []
        inherited_state = inherited_state or {}
        current_heading = inherited_state.get("heading", "")
        current_chapter_number = int(inherited_state.get("chapter_number") or 0)
        current_chapter_title = inherited_state.get("chapter_title", "")
        current_section_number = inherited_state.get("section_number", "")
        current_section_title = inherited_state.get("section_title", "")
        current_parts = []

        def flush():
            if not current_parts:
                return
            unit_text = self.clean_text("\n".join(current_parts))
            if not unit_text:
                return
            detected = self.detect_heading(unit_text)
            heading = current_heading or detected
            section_path = self.make_section_path(
                current_chapter_number,
                current_chapter_title,
                current_section_number,
                current_section_title,
                heading
            )
            units.append({
                "text": unit_text,
                "page_number": page_number,
                "slide_number": slide_number,
                "heading": heading,
                "section_path": section_path,
                "detected_title": detected or heading,
                "chapter_number": current_chapter_number,
                "chapter_title": current_chapter_title,
                "section_number": current_section_number,
                "section_title": current_section_title,
                "content_zone": self.detect_content_zone(unit_text, page_number, heading)
            })

        for line in lines:
            if not line:
                if current_parts:
                    current_parts.append("")
                continue
            chapter = self.parse_chapter_heading(line)
            if (
                chapter
                and not chapter.get("title")
                and current_chapter_number
                and int(chapter.get("number") or 0) != current_chapter_number
            ):
                chapter = None
            section = self.parse_section_heading(line)
            heading = self.detect_heading_line(line)
            if (chapter or section or heading) and current_parts:
                flush()
                current_parts = []
            if chapter:
                current_chapter_number = chapter["number"]
                current_chapter_title = chapter.get("title", "") or current_chapter_title
                current_section_number = ""
                current_section_title = ""
                current_heading = self.make_section_path(current_chapter_number, current_chapter_title)
            elif section:
                current_section_number = section["number"]
                current_section_title = section["title"]
                major = int(current_section_number.split(".", 1)[0])
                if major and current_chapter_number != major:
                    current_chapter_number = major
                    current_chapter_title = {
                        1: "Introduction",
                        2: "Overview"
                    }.get(major, current_chapter_title if current_chapter_number == major else "")
                current_heading = self.make_section_path(
                    current_chapter_number,
                    current_chapter_title,
                    current_section_number,
                    current_section_title
                )
            if heading:
                current_heading = heading
            current_parts.append(line)

        flush()
        return units, {
            "heading": current_heading,
            "chapter_number": current_chapter_number,
            "chapter_title": current_chapter_title,
            "section_number": current_section_number,
            "section_title": current_section_title
        }

    def extract_units_from_pdf(self, file_path):
        units = []
        current_state = {}
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page_number, page in enumerate(reader.pages, 1):
                text = self.clean_text(page.extract_text() or "")
                if not text:
                    continue
                page_chapter = self.detect_page_chapter(text)
                if page_chapter:
                    current_state = {
                        "heading": self.make_section_path(page_chapter["number"], page_chapter.get("title", "")),
                        "chapter_number": page_chapter["number"],
                        "chapter_title": page_chapter.get("title", ""),
                        "section_number": "",
                        "section_title": ""
                    }
                page_units, current_state = self.split_text_into_units(
                    text,
                    page_number=page_number,
                    inherited_state=current_state
                )
                units.extend(page_units)
        return units

    def extract_units_from_docx(self, file_path):
        doc = DocxDocument(file_path)
        units = []
        current_heading = ""
        current_parts = []

        def flush():
            if not current_parts:
                return
            text = self.clean_text("\n".join(current_parts))
            if not text:
                return
            detected = self.detect_heading(text)
            heading = current_heading or detected
            units.append({
                "text": text,
                "page_number": None,
                "slide_number": None,
                "heading": heading,
                "section_path": heading,
                "detected_title": detected or heading,
                "chapter_number": 0,
                "chapter_title": "",
                "section_number": "",
                "section_title": "",
                "content_zone": self.detect_content_zone(text, None, heading)
            })

        for para in doc.paragraphs:
            text = self.clean_text(para.text)
            if not text:
                continue
            style_name = (para.style.name if para.style else "").lower()
            looks_like_heading = style_name.startswith("heading") or bool(self.detect_heading_line(text))
            if looks_like_heading and current_parts:
                flush()
                current_parts = []
            if looks_like_heading:
                current_heading = text[:180]
            current_parts.append(text)

        flush()
        return units

    def extract_units_from_pptx(self, file_path):
        prs = Presentation(file_path)
        units = []
        for slide_number, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            text = self.clean_text("\n".join(slide_text))
            if not text:
                continue
            slide_units, _ = self.split_text_into_units(text, slide_number=slide_number)
            units.extend(slide_units)
        return units

    def merge_short_chunks(self, split_texts, min_chars=180):
        merged = []
        pending = ""
        for text in split_texts:
            clean = self.clean_text(text)
            if not clean:
                continue
            pending = self.clean_text(f"{pending}\n{clean}") if pending else clean
            if len(pending) >= min_chars:
                merged.append(pending)
                pending = ""

        if pending:
            if merged and len(pending) < min_chars:
                merged[-1] = self.clean_text(f"{merged[-1]}\n{pending}")
            else:
                merged.append(pending)
        return merged

    def split_units(self, units):
        chunks = []
        for unit in units:
            split_texts = self.merge_short_chunks(self.text_splitter.split_text(unit["text"]))
            for local_index, chunk_text in enumerate(split_texts):
                chunk_text = self.clean_text(chunk_text)
                if len(chunk_text) < 20:
                    continue
                chunks.append({
                    "text": chunk_text,
                    "page_number": unit.get("page_number"),
                    "slide_number": unit.get("slide_number"),
                    "heading": unit.get("heading") or "",
                    "section_path": unit.get("section_path") or unit.get("heading") or "",
                    "detected_title": unit.get("detected_title") or unit.get("heading") or "",
                    "chapter_number": int(unit.get("chapter_number") or 0),
                    "chapter_title": unit.get("chapter_title") or "",
                    "section_number": unit.get("section_number") or "",
                    "section_title": unit.get("section_title") or "",
                    "content_zone": unit.get("content_zone") or "body",
                    "local_index": local_index
                })
        return chunks

    def process_file(self, file_path):
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            units = self.extract_units_from_pdf(file_path)
        elif ext == ".docx":
            units = self.extract_units_from_docx(file_path)
        elif ext in [".pptx", ".ppt"]:
            units = self.extract_units_from_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file extension: {ext}")

        extracted_chars = sum(len(unit["text"]) for unit in units)
        if extracted_chars < 50:
            print(f"  Warning: Very little text extracted ({extracted_chars} chars)", flush=True)
            return []

        chunks = self.split_units(units)
        print(
            f"  Extracted {extracted_chars} chars -> {len(units)} sections -> {len(chunks)} structured chunks "
            f"(size={self.chunk_size}, overlap={self.chunk_overlap})",
            flush=True
        )
        return chunks
